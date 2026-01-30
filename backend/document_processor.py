#Agentic RAG System/backend/document_processor.py
import os
import re
from typing import List, Tuple
import pandas as pd
from config import Config

class DocumentProcessor:
    def __init__(self):
        self.chunk_size = 500
        self.chunk_overlap = 50
    
    def process_document(self, filepath: str, filename: str) -> Tuple[List[str], List[dict]]:
        """Process document and return chunks with metadata"""
        ext = os.path.splitext(filename)[1].lower()
        
        if ext == '.pdf':
            text = self._extract_pdf(filepath)
        elif ext == '.docx':
            text = self._extract_docx(filepath)
        elif ext == '.txt':
            text = self._extract_text(filepath)
        elif ext in ['.xlsx', '.xls']:
            text = self._extract_excel(filepath)
        elif ext == '.pptx':
            text = self._extract_pptx(filepath)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        
        # Clean text
        text = self._clean_text(text)
        
        # Chunk text
        chunks = self._chunk_text(text)
        
        # Create metadata
        metadata = []
        for i, chunk in enumerate(chunks):
            metadata.append({
                "source": filename,
                "chunk_id": i,
                "total_chunks": len(chunks),
                "file_type": ext[1:],  # Remove dot
                "chunk_length": len(chunk)
            })
        
        return chunks, metadata
    
    def _extract_pdf(self, filepath: str) -> str:
        """Extract text from PDF"""
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += f"--- Page {page.page_number} ---\n{page_text}\n\n"
            return text
        except ImportError:
            # Fallback to PyPDF2
            import PyPDF2
            text = ""
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
            return text
    
    def _extract_docx(self, filepath: str) -> str:
        """Extract text from DOCX"""
        from docx import Document
        doc = Document(filepath)
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        return text
    
    def _extract_text(self, filepath: str) -> str:
        """Extract text from TXT"""
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    def _extract_excel(self, filepath: str) -> str:
        """Extract text from Excel"""
        text = ""
        try:
            excel_data = pd.read_excel(filepath, sheet_name=None)
            for sheet_name, df in excel_data.items():
                text += f"--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n\n"
        except Exception as e:
            text = f"Error reading Excel: {str(e)}"
        return text
    
    def _extract_pptx(self, filepath: str) -> str:
        """Extract text from PowerPoint"""
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            text += "\n"
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove page/slide markers
        text = re.sub(r'--- (Page|Slide|Sheet) \d+ ---', '', text)
        # Remove special characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
        return text.strip()
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks"""
        if not text:
            return []
        
        # Split by paragraphs first
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = ""
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            if len(current_chunk) + len(para) > self.chunk_size and current_chunk:
                chunks.append(current_chunk)
                # Keep overlap
                words = current_chunk.split()
                overlap_words = words[-min(len(words), self.chunk_overlap):]
                current_chunk = " ".join(overlap_words) + " " + para if overlap_words else para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
        
        if current_chunk:
            chunks.append(current_chunk)
        
        # If chunks are too large, split further
        final_chunks = []
        for chunk in chunks:
            if len(chunk) > self.chunk_size * 1.5:
                # Split by sentences
                sentences = re.split(r'(?<=[.!?])\s+', chunk)
                temp_chunk = ""
                for sentence in sentences:
                    if len(temp_chunk) + len(sentence) > self.chunk_size and temp_chunk:
                        final_chunks.append(temp_chunk)
                        temp_chunk = sentence
                    else:
                        if temp_chunk:
                            temp_chunk += " " + sentence
                        else:
                            temp_chunk = sentence
                if temp_chunk:
                    final_chunks.append(temp_chunk)
            else:
                final_chunks.append(chunk)
        
        return final_chunks